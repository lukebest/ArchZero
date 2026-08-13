"""Patent disclosure: benefits come from funnel metrics, and the deck says so.

The pptx round-trip needs the optional extra (``uv sync --extra patent``) and
skips cleanly without it.
"""

from __future__ import annotations

import json

import pytest

from archzero.llm.fake import FakeLLM
from archzero.models import (
    Candidate,
    EvidenceLevel,
    Tier,
    TierResult,
    Verdict,
)
from archzero.patent.disclosure import (
    build_disclosure,
    collect_benefits,
    disclosure_markdown,
)
from archzero.patent.models import SECTION_TITLES, PatentDisclosure
from archzero.patent.prior_art import PriorArtResult

pytestmark = pytest.mark.patent


@pytest.fixture
def pptx_mod():
    return pytest.importorskip("pptx", reason="needs: uv sync --extra patent")


@pytest.fixture
def measured_candidate() -> Candidate:
    """A candidate that made it to Tier3 with real numbers attached."""
    cand = Candidate(
        problem_id="pp-demo-cache",
        title="死块过滤预取",
        mechanism="用死块预测过滤 L2 预取请求，仅对预测存活的块下发预取。",
        family="prefetch",
        metrics={
            "t2_miss_reduction": 0.21,
            "t3_miss_reduction": 0.184,
            "t3_bw_delta_frac": 0.031,
            "t3_area_mm2": 0.24,
            "t3_magic_gap": 1.14,
            "diverge_lens": "queueing_theory",
            "diverge_domain": "tcp_congestion",
            "diverge_mode": "lateral",
        },
    )
    cand.tier_history = [
        TierResult(
            tier=Tier.T2,
            verdict=Verdict.PASS,
            score=0.8,
            summary="解析模型达标",
            evidence=EvidenceLevel.ANALYTIC,
            model_id="fake-model",
        ),
        TierResult(
            tier=Tier.T3,
            verdict=Verdict.PASS,
            score=0.75,
            summary="小规模仿真达标",
            evidence=EvidenceLevel.SIM,
            model_id="fake-model",
        ),
    ]
    return cand


@pytest.fixture
def draft_llm() -> FakeLLM:
    return FakeLLM(
        responses={
            "prior_art": json.dumps({"queries": ["dead block prefetch filtering"]}),
            "patent_draft": json.dumps(
                {
                    "background": "在 LLM 推理场景下 L2 缺失率主导端到端时延。",
                    "existing_tech": "主流做法是无差别流预取，带宽浪费严重。",
                    "solution": "在预取队列前插入死块预测过滤器。",
                    "steps": ["采样访问历史", "训练死块预测表", "过滤预取请求"],
                    "points": [
                        {
                            "title": "一种基于死块预测的预取过滤方法",
                            "essential_features": ["死块预测表", "预取请求过滤门"],
                            "depends_on": None,
                        },
                        {
                            "title": "带宽预算反馈的过滤阈值调整",
                            "essential_features": ["带宽计量器", "阈值反馈回路"],
                            "depends_on": 1,
                        },
                    ],
                    "statements": {
                        "t3_miss_reduction": "显著降低二级缓存缺失率，满足验收门限。",
                        "t3_bw_delta_frac": "DRAM 带宽开销控制在验收门限之内。",
                    },
                    "qualitative": ["不改变现有缓存一致性协议，改动面小。"],
                },
                ensure_ascii=False,
            ),
        }
    )


@pytest.fixture
def force_offline(monkeypatch):
    from archzero.patent import prior_art as pa

    async def _boom(*args, **kwargs):
        raise RuntimeError("connection refused")

    monkeypatch.setattr(pa, "search_arxiv", _boom)
    monkeypatch.setattr(pa, "search_semantic_scholar", _boom)


def test_benefits_come_from_metrics_only(measured_candidate, demo_problem):
    benefits = collect_benefits(measured_candidate, demo_problem)
    by_key = {b.metric_key: b for b in benefits}

    assert by_key["t3_miss_reduction"].metric_value == 0.184
    assert by_key["t3_miss_reduction"].display_value == "18.4%"
    assert by_key["t3_miss_reduction"].evidence_level == "sim"
    assert by_key["t3_miss_reduction"].tier == "tier3"
    # Sim evidence supersedes the Tier2 analytic estimate for the same quantity.
    assert "t2_miss_reduction" not in by_key


def test_benefits_are_checked_against_acceptance_thresholds(
    measured_candidate, demo_problem
):
    by_key = {b.metric_key: b for b in collect_benefits(measured_candidate, demo_problem)}

    assert by_key["t3_miss_reduction"].meets_threshold is True
    assert by_key["t3_bw_delta_frac"].meets_threshold is True
    assert "门限" in by_key["t3_miss_reduction"].threshold


def test_unmeasured_candidate_yields_no_quantified_benefits(demo_problem):
    bare = Candidate(problem_id="pp", title="未评估机制", mechanism="仅有构想。")
    assert collect_benefits(bare, demo_problem) == []


@pytest.mark.asyncio
async def test_disclosure_has_all_six_sections(
    tmp_cfg, measured_candidate, demo_problem, draft_llm, force_offline
):
    disc = await build_disclosure(
        tmp_cfg, measured_candidate, problem=demo_problem, llm=draft_llm
    )

    assert disc.background
    assert disc.existing_tech
    assert disc.technical_solution
    assert len(disc.solution_steps) == 3
    assert len(disc.protection_points) == 2
    assert disc.protection_points[0].depends_on is None
    assert disc.protection_points[1].depends_on == 1
    assert disc.quantified_benefits
    assert disc.prior_art.retrieval_status == "offline"

    md = disclosure_markdown(disc)
    for title in SECTION_TITLES:
        assert title in md
    assert "18.4%" in md
    assert disc.evidence_level == "sim"


@pytest.mark.asyncio
async def test_llm_wording_cannot_overwrite_the_numbers(
    tmp_cfg, measured_candidate, demo_problem, force_offline
):
    """A model claiming 90% must not change what the deck reports."""
    liar = FakeLLM(
        responses={
            "prior_art": json.dumps({"queries": ["x"]}),
            "patent_draft": json.dumps(
                {
                    "statements": {"t3_miss_reduction": "缺失率下降 90%，性能翻倍。"},
                    "points": [],
                },
                ensure_ascii=False,
            ),
        }
    )
    disc = await build_disclosure(
        tmp_cfg, measured_candidate, problem=demo_problem, llm=liar
    )
    benefit = next(b for b in disc.benefits if b.metric_key == "t3_miss_reduction")

    assert benefit.metric_value == 0.184
    assert benefit.display_value == "18.4%"
    assert "90" not in benefit.statement
    assert any("自造数字" in w for w in disc.warnings)
    assert "90%" not in disclosure_markdown(disc)


def _deck_text(path) -> str:
    from pptx import Presentation

    chunks: list[str] = []
    for slide in Presentation(str(path)).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    chunks.extend(cell.text for cell in row.cells)
    return "\n".join(chunks)


@pytest.mark.asyncio
async def test_deck_round_trips_all_six_sections_and_numbers(
    tmp_cfg, tmp_path, measured_candidate, demo_problem, draft_llm, force_offline, pptx_mod
):
    from archzero.patent.pptx_render import render_deck

    disc = await build_disclosure(
        tmp_cfg, measured_candidate, problem=demo_problem, llm=draft_llm
    )
    out = render_deck(tmp_cfg, disc, tmp_path / "deck.pptx")
    assert out.is_file()

    text = _deck_text(out)
    for title in SECTION_TITLES:
        assert title in text, f"missing section: {title}"

    # Numbers on the slide must equal the candidate's measured metrics.
    for benefit in disc.quantified_benefits:
        assert benefit.display_value in text
    assert "18.4%" in text
    assert "90%" not in text

    # Honesty markers for the degraded retrieval and the provenance appendix.
    assert "offline" in text
    assert "待人工核实" in text
    assert "证据溯源" in text
    assert "tcp_congestion" in text


@pytest.mark.asyncio
async def test_deck_is_16_by_9(tmp_cfg, tmp_path, pptx_mod):
    from pptx import Presentation
    from pptx.util import Inches

    from archzero.patent.pptx_render import render_deck

    disc = PatentDisclosure(
        candidate_id="cand-1",
        title="空壳方案",
        prior_art=PriorArtResult(),
    )
    out = render_deck(tmp_cfg, disc, tmp_path / "empty.pptx")
    prs = Presentation(str(out))

    assert prs.slide_width == Inches(13.333)
    assert prs.slide_height == Inches(7.5)
    assert len(prs.slides) >= 7


def test_long_prose_spills_onto_continuation_slides(tmp_cfg, tmp_path, pptx_mod):
    from pptx import Presentation

    from archzero.patent.pptx_render import render_deck

    disc = PatentDisclosure(
        candidate_id="cand-2",
        title="长文方案",
        background="背" * 4000,
    )
    out = render_deck(tmp_cfg, disc, tmp_path / "long.pptx")
    text = "\n".join(
        shape.text_frame.text
        for slide in Presentation(str(out)).slides
        for shape in slide.shapes
        if shape.has_text_frame
    )

    assert f"{SECTION_TITLES[0]}（1/" in text
    assert f"{SECTION_TITLES[0]}（2/" in text


def test_east_asian_font_is_written(tmp_cfg, tmp_path, pptx_mod):
    from pptx import Presentation
    from pptx.oxml.ns import qn

    from archzero.patent.pptx_render import render_deck

    tmp_cfg.patent.ea_font = "Source Han Sans SC"
    disc = PatentDisclosure(candidate_id="cand-3", title="中文字体检查")
    out = render_deck(tmp_cfg, disc, tmp_path / "font.pptx")

    ea_faces = {
        node.get("typeface")
        for slide in Presentation(str(out)).slides
        for shape in slide.shapes
        if shape.has_text_frame
        for node in shape.text_frame._txBody.iter(qn("a:ea"))
    }
    assert ea_faces == {"Source Han Sans SC"}


def test_missing_dependency_raises_actionable_error(tmp_cfg, tmp_path, monkeypatch):
    import builtins

    from archzero.patent import PatentDepsMissing
    from archzero.patent.pptx_render import render_deck

    real_import = builtins.__import__

    def _blocked(name, *args, **kwargs):
        if name == "pptx" or name.startswith("pptx."):
            raise ImportError("No module named 'pptx'")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", _blocked)
    disc = PatentDisclosure(candidate_id="cand-4", title="无依赖")

    with pytest.raises(PatentDepsMissing) as exc:
        render_deck(tmp_cfg, disc, tmp_path / "nope.pptx")
    assert "uv sync --extra patent" in str(exc.value)
    assert "--md-only" in str(exc.value)
