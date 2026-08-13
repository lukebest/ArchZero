"""Render a PatentDisclosure into a 16:9 review deck.

Optional module: needs ``uv sync --extra patent``. ``pptx`` is imported inside
the functions so importing this file is always safe.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from archzero.config import FactoryConfig
from archzero.patent import PatentDepsMissing
from archzero.patent.models import SECTION_TITLES, PatentDisclosure

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
BODY_LINES_PER_SLIDE = 14
TABLE_ROWS_PER_SLIDE = 8

_THREAT_LABEL = {"high": "高", "medium": "中", "low": "低"}


def _require_pptx() -> Any:
    try:
        import pptx
    except ImportError as exc:  # pragma: no cover - depends on install profile
        raise PatentDepsMissing("python-pptx") from exc
    return pptx


def _set_font(run: Any, cfg: FactoryConfig, size_pt: int, *, bold: bool = False) -> None:
    """Set latin + East Asian typeface.

    ``font.name`` only writes ``<a:latin>``. Without an explicit ``<a:ea>`` the
    Chinese glyphs fall back to whatever the renderer picks, which usually
    means mismatched weights and broken line spacing.
    """
    from pptx.oxml.ns import qn
    from pptx.util import Pt

    font = run.font
    font.size = Pt(size_pt)
    font.bold = bold
    font.name = cfg.patent.latin_font

    rpr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        for existing in rpr.findall(qn(tag)):
            rpr.remove(existing)
        node = rpr.makeelement(qn(tag), {"typeface": cfg.patent.ea_font})
        rpr.append(node)


def _blank(prs: Any) -> Any:
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_title(slide: Any, cfg: FactoryConfig, text: str, *, subtitle: str = "") -> Any:
    from pptx.util import Inches

    box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.35), Inches(SLIDE_W_IN - 1.2), Inches(1.0)
    )
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    _set_font(run, cfg, 28, bold=True)
    if subtitle:
        sub = tf.add_paragraph()
        srun = sub.add_run()
        srun.text = subtitle
        _set_font(srun, cfg, 13)
    return box


def _add_bullets(
    slide: Any,
    cfg: FactoryConfig,
    lines: list[tuple[str, int]],
    *,
    top_in: float = 1.5,
    size_pt: int = 16,
) -> None:
    """``lines`` is (text, indent_level)."""
    from pptx.util import Inches

    box = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(top_in),
        Inches(SLIDE_W_IN - 1.4),
        Inches(SLIDE_H_IN - top_in - 0.6),
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = min(max(level, 0), 4)
        run = para.add_run()
        run.text = text
        _set_font(run, cfg, size_pt if level == 0 else size_pt - 2)


def _wrap(text: str, width: int = 46) -> list[str]:
    """Split prose into slide-sized lines. CJK has no spaces, so cut by width."""
    out: list[str] = []
    for raw in (text or "").split("\n"):
        raw = raw.strip()
        if not raw:
            continue
        while len(raw) > width:
            out.append(raw[:width])
            raw = raw[width:]
        if raw:
            out.append(raw)
    return out


def _prose_slides(
    prs: Any,
    cfg: FactoryConfig,
    title: str,
    text: str,
    *,
    empty_hint: str = "（本节未生成，需人工补充）",
) -> None:
    lines = _wrap(text) or [empty_hint]
    chunks = [
        lines[i : i + BODY_LINES_PER_SLIDE]
        for i in range(0, len(lines), BODY_LINES_PER_SLIDE)
    ]
    for idx, chunk in enumerate(chunks):
        slide = _blank(prs)
        suffix = "" if len(chunks) == 1 else f"（{idx + 1}/{len(chunks)}）"
        _add_title(slide, cfg, f"{title}{suffix}")
        _add_bullets(slide, cfg, [(line, 0) for line in chunk])


def _table_slides(
    prs: Any,
    cfg: FactoryConfig,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    empty_hint: str,
) -> None:
    from pptx.util import Inches

    if not rows:
        slide = _blank(prs)
        _add_title(slide, cfg, title)
        _add_bullets(slide, cfg, [(empty_hint, 0)])
        return

    chunks = [
        rows[i : i + TABLE_ROWS_PER_SLIDE]
        for i in range(0, len(rows), TABLE_ROWS_PER_SLIDE)
    ]
    for idx, chunk in enumerate(chunks):
        slide = _blank(prs)
        suffix = "" if len(chunks) == 1 else f"（{idx + 1}/{len(chunks)}）"
        _add_title(slide, cfg, f"{title}{suffix}")
        shape = slide.shapes.add_table(
            len(chunk) + 1,
            len(headers),
            Inches(0.6),
            Inches(1.5),
            Inches(SLIDE_W_IN - 1.2),
            Inches(0.4 * (len(chunk) + 1)),
        )
        table = shape.table
        for c, head in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = ""
            run = cell.text_frame.paragraphs[0].add_run()
            run.text = head
            _set_font(run, cfg, 12, bold=True)
        for r, row in enumerate(chunk, start=1):
            for c, value in enumerate(row):
                cell = table.cell(r, c)
                cell.text = ""
                para = cell.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = value
                _set_font(run, cfg, 11)


def render_deck(
    cfg: FactoryConfig,
    disc: PatentDisclosure,
    out_path: Path,
) -> Path:
    """Write the six-section review deck. Raises PatentDepsMissing without the extra."""
    pptx = _require_pptx()
    from pptx.util import Inches

    prs = pptx.Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Cover
    cover = _blank(prs)
    _add_title(
        cover,
        cfg,
        f"专利评审材料：{disc.title}",
        subtitle=f"所属问题：{disc.problem_title or '（未关联问题包）'}",
    )
    _add_bullets(
        cover,
        cfg,
        [
            (f"候选 ID：{disc.candidate_id}", 0),
            (f"机制族：{disc.family}", 0),
            (f"最强证据等级：{disc.evidence_level}", 0),
            (f"检索状态：{disc.prior_art.retrieval_status}", 0),
            (f"生成时间：{disc.created_at.strftime('%Y-%m-%d %H:%M UTC')}", 0),
        ],
        top_in=2.4,
        size_pt=18,
    )
    if not disc.prior_art.verified or disc.warnings:
        _add_bullets(
            cover,
            cfg,
            [("注意：本材料含未经人工核实的内容，详见末页告警", 0)],
            top_in=5.6,
            size_pt=14,
        )

    # 1 / 2 / 3
    _prose_slides(prs, cfg, SECTION_TITLES[0], disc.background)
    _prose_slides(prs, cfg, SECTION_TITLES[1], disc.existing_tech)
    _prose_slides(prs, cfg, SECTION_TITLES[2], disc.technical_solution)
    if disc.solution_steps:
        slide = _blank(prs)
        _add_title(slide, cfg, f"{SECTION_TITLES[2]} · 实施步骤")
        _add_bullets(
            slide,
            cfg,
            [(f"{i}. {s}", 0) for i, s in enumerate(disc.solution_steps, start=1)],
        )

    # 4 保护点
    if disc.protection_points:
        for point in disc.protection_points:
            slide = _blank(prs)
            _add_title(
                slide,
                cfg,
                f"{SECTION_TITLES[3]} · 保护点 {point.index}",
                subtitle=point.kind,
            )
            lines: list[tuple[str, int]] = [(point.title, 0)]
            lines += [("必要技术特征：", 0)]
            lines += [(f, 1) for f in point.essential_features]
            _add_bullets(slide, cfg, lines)
    else:
        slide = _blank(prs)
        _add_title(slide, cfg, SECTION_TITLES[3])
        _add_bullets(slide, cfg, [("（未生成保护点，需人工补充）", 0)])

    # 5 有益效果
    quantified = [b for b in disc.benefits if b.quantified]
    _table_slides(
        prs,
        cfg,
        SECTION_TITLES[4],
        ["效果", "实测值", "验收门限", "达标", "证据", "来源"],
        [
            [
                b.statement,
                b.display_value,
                b.threshold or "-",
                "-" if b.meets_threshold is None else ("是" if b.meets_threshold else "否"),
                b.evidence_level,
                b.tier or "-",
            ]
            for b in quantified
        ],
        empty_hint=disc.benefits_note or "（无实测指标，本节需补充仿真数据）",
    )
    qualitative = [b for b in disc.benefits if not b.quantified]
    if qualitative:
        slide = _blank(prs)
        _add_title(slide, cfg, f"{SECTION_TITLES[4]} · 定性效果")
        _add_bullets(slide, cfg, [(b.statement, 0) for b in qualitative])

    # 6 检索与对比
    pa = disc.prior_art
    _table_slides(
        prs,
        cfg,
        SECTION_TITLES[5],
        ["对比文献", "相同点", "不同点", "区别性技术特征", "威胁"],
        [
            [
                c.title,
                c.same_points or "-",
                c.diff_points or "-",
                c.distinguishing_features or "-",
                _THREAT_LABEL.get(c.threat, c.threat),
            ]
            for c in pa.comparisons
        ],
        empty_hint="（未完成逐条对比）",
    )

    slide = _blank(prs)
    _add_title(slide, cfg, f"{SECTION_TITLES[5]} · 检索说明")
    search_lines: list[tuple[str, int]] = [
        (pa.caveat(), 0),
        (f"检索状态：{pa.retrieval_status}", 0),
        (f"命中文献：{len(pa.hits)} 篇", 0),
    ]
    if pa.queries:
        search_lines.append(("使用的检索式：", 0))
        search_lines += [(q, 1) for q in pa.queries]
    sug = pa.patent_suggestion
    if sug.ipc or sug.queries_cn or sug.queries_en:
        search_lines.append(("专利库检索建议（需人工在内部专利库执行）：", 0))
        if sug.ipc:
            search_lines.append((f"IPC/CPC：{', '.join(sug.ipc)}", 1))
        search_lines += [(q, 1) for q in sug.queries_cn]
        search_lines += [(q, 1) for q in sug.queries_en]
    _add_bullets(slide, cfg, search_lines, size_pt=14)

    # Appendix: provenance
    slide = _blank(prs)
    _add_title(slide, cfg, "附录：证据溯源", subtitle="由 ArchZero 漏斗自动记录")
    prov_lines: list[tuple[str, int]] = [
        (f"候选 ID：{disc.candidate_id}", 0),
        (f"机制族：{disc.family}", 0),
    ]
    dv = disc.provenance.get("diverge_domain")
    if dv:
        prov_lines.append(
            (
                f"跨域来源：{dv} / 透镜：{disc.provenance.get('diverge_lens')} "
                f"/ 模式：{disc.provenance.get('diverge_mode')}",
                0,
            )
        )
    prov_lines.append(("Tier 历史：", 0))
    for row in disc.provenance.get("tier_history", []):
        prov_lines.append(
            (
                f"{row['tier']} · {row['verdict']} · evidence={row['evidence']} "
                f"· model={row.get('model_id') or '-'}",
                1,
            )
        )
    if disc.warnings:
        prov_lines.append(("生成告警：", 0))
        prov_lines += [(w, 1) for w in disc.warnings]
    if pa.notes:
        prov_lines.append(("检索备注：", 0))
        prov_lines += [(n, 1) for n in pa.notes]
    _add_bullets(slide, cfg, prov_lines, size_pt=13)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
